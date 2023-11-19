from pptx import Presentation
import boto3
import io
from botocore.exceptions import ClientError
import os

s3 = boto3.client("s3")
OUTPUT_BUCKET = os.environ["OUTPUT_BUCKET"]


def handler(event, context):
    bucket_name = event["Records"][0]["s3"]["bucket"]["name"]
    object_key = event["Records"][0]["s3"]["object"]["key"].replace("+", " ")

    try:
        response = s3.get_object(Bucket=bucket_name, Key=object_key)
        data = response["Body"].read()
        remote_file_bytes = io.BytesIO(data)

        prs = Presentation(remote_file_bytes)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text + " "
                        text += "\n"

        response = s3.put_object(
            Body=text, Bucket=OUTPUT_BUCKET, Key=object_key + ".txt"
        )

    except ClientError as e:
        print(f"Error parsing the word file : {e}")
